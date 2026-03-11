#!/usr/bin/env python3
"""
Build PPTX directly from Marp slides.md — faithful dark tech theme.

Adapted from DSC 190 presentation make_pptx.py.
No Pandoc. Every element is constructed with python-pptx: positioned text
boxes, native tables, styled shapes.  The result mirrors the original Marp
CSS (colors, fonts, spacing, decorations) while keeping every element
individually editable in PowerPoint / Keynote / Google Slides.

Requirements: python-pptx, lxml
Optional:     Inter + Fira Code fonts (falls back to system fonts)
"""

import re
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml.etree import SubElement

DIR = Path(__file__).parent

# ── Palette  (from slides.md :root) ────────────────────────────
C_BG      = RGBColor(0x0D, 0x11, 0x17)
C_FG      = RGBColor(0xC9, 0xD1, 0xD9)
C_HEAD    = RGBColor(0x58, 0xA6, 0xFF)
C_ACCENT  = RGBColor(0x7E, 0xE7, 0x87)
C_ORANGE  = RGBColor(0xFF, 0xA6, 0x57)
C_PURPLE  = RGBColor(0xD2, 0xA8, 0xFF)
C_RED     = RGBColor(0xF8, 0x51, 0x49)
C_CODE_BG = RGBColor(0x16, 0x1B, 0x22)
C_BORDER  = RGBColor(0x30, 0x36, 0x3D)
C_GRAY    = RGBColor(0x8B, 0x94, 0x9E)
C_DIM     = RGBColor(0x48, 0x4F, 0x58)

FB = "Inter"          # body font
FC = "Fira Code"      # heading / code font

# Slide geometry  (16 : 9)
SW, SH = Inches(13.333), Inches(7.5)
MX     = Inches(0.78)               # left / right margin
CW     = SW - 2 * MX                # content width
STRIPE = Emu(48000)                  # green left-border width

# Font sizes
Z = dict(
    h1=Pt(38), h2=Pt(28), h3=Pt(20),
    body=Pt(16), small=Pt(13),
    th=Pt(14), td=Pt(13),
    footer=Pt(10),
    sec_title=Pt(34), sec_sub=Pt(18),
    lead_body=Pt(18),
    code=Pt(12),
)

# Layout positions (EMU)
TITLE_Y  = Inches(0.50)
TITLE_H  = Inches(0.72)
BORDER_Y = TITLE_Y + TITLE_H + Emu(36000)
BODY_Y   = BORDER_Y + Emu(54000)
BODY_MAX = Inches(6.65) - BODY_Y     # max body height

FOOTER_TEXT = "LOCALE | DSC 190 | Prof. Biwei Huang"


# ═══════════════════════════════════════════════════════════════
#  PARSING
# ═══════════════════════════════════════════════════════════════

def parse_inline(text):
    """Return [(text, style)] where style ∈ {n, b, i, c}."""
    out, last = [], 0
    for m in re.finditer(r"\*\*(.*?)\*\*|\*([^*]+?)\*|`([^`]+?)`", text):
        if m.start() > last:
            out.append((text[last : m.start()], "n"))
        if   m.group(1) is not None: out.append((m.group(1), "b"))
        elif m.group(2) is not None: out.append((m.group(2), "i"))
        elif m.group(3) is not None: out.append((m.group(3), "c"))
        last = m.end()
    if last < len(text):
        out.append((text[last:], "n"))
    return out or [("", "n")]


def parse_body(md):
    """Markdown body → list of typed content blocks."""
    blocks, buf = [], None
    in_code = False
    code_lines = []

    def flush():
        nonlocal buf
        if buf:
            blocks.append(buf)
            buf = None

    for line in md.split("\n"):
        s = line.strip()

        # Code block toggle
        if s.startswith("```"):
            if in_code:
                # End code block
                in_code = False
                flush()
                blocks.append({"type": "code", "text": "\n".join(code_lines)})
                code_lines = []
                continue
            else:
                # Start code block
                in_code = True
                flush()
                continue

        if in_code:
            code_lines.append(line.rstrip())
            continue

        if not s:
            flush()
            continue

        # Blockquote
        if s.startswith("> "):
            flush()
            blocks.append({"type": "quote", "text": s[2:]})
            continue

        if s.startswith("### "):
            flush()
            blocks.append({"type": "h3", "text": s[4:]})
            continue

        ml = re.match(r"^(\s*)(?:[-*]|(\d+)[.)])\s+(.+)", line)
        if ml:
            indent, num, text = len(ml.group(1)), ml.group(2), ml.group(3)
            if not buf or buf["type"] != "list":
                flush()
                buf = {"type": "list", "items": []}
            buf["items"].append({"text": text, "level": min(indent // 2, 2), "num": num})
            continue

        if s.startswith("|") and s.endswith("|"):
            if re.match(r"^\|[\s\-:|]+\|$", s):
                continue
            cells = [c.strip() for c in s.split("|")[1:-1]]
            if not buf or buf["type"] != "table":
                flush()
                buf = {"type": "table", "rows": []}
            buf["rows"].append(cells)
            continue

        if buf and buf["type"] == "para":
            buf["text"] += " " + s
        else:
            flush()
            buf = {"type": "para", "text": s}

    # Close any open code block
    if in_code and code_lines:
        flush()
        blocks.append({"type": "code", "text": "\n".join(code_lines)})

    flush()
    return blocks


def parse_slides(src):
    raw = src.read_text(encoding="utf-8")
    raw = re.sub(r"^---\n.*?^---\n", "", raw, count=1, flags=re.DOTALL | re.MULTILINE)
    raw = re.sub(r"<style>.*?</style>\s*", "", raw, flags=re.DOTALL)

    slides = []
    for i, part in enumerate(re.split(r"\n---\n", raw)):
        part = part.strip()
        if not part:
            continue

        is_title = "_class: title" in part
        is_sec   = "_class: section-break" in part
        is_lead  = "_class: lead" in part
        part = re.sub(r"<!--.*?-->", "", part, flags=re.DOTALL).strip()
        # Remove <span> tags but keep content
        part = re.sub(r"<span[^>]*>(.*?)</span>", r"\1", part, flags=re.DOTALL)
        if not part:
            continue

        lines = part.split("\n")
        titles, body_start = [], 0
        for j, ln in enumerate(lines):
            if re.match(r"^#{1,2}\s", ln.strip()):
                titles.append(re.sub(r"^#{1,2}\s*", "", ln.strip()))
                body_start = j + 1
            elif ln.strip() == "":
                body_start = j + 1
            else:
                break

        title = "\n".join(titles)
        body  = "\n".join(lines[body_start:]).strip()
        stype = (
            "title"   if is_title or (i == 0 and is_lead) else
            "section" if is_sec else
            "lead"    if is_lead else
            "content"
        )
        slides.append(dict(type=stype, title=title, body=body,
                           blocks=parse_body(body) if body else []))
    return slides


# ═══════════════════════════════════════════════════════════════
#  LOW-LEVEL HELPERS
# ═══════════════════════════════════════════════════════════════

def _r(para, txt, font, sz, clr, bold=False, italic=False):
    """Append a styled run."""
    r = para.add_run()
    r.text = txt
    r.font.name  = font
    r.font.size  = sz
    r.font.color.rgb = clr
    if bold:   r.font.bold   = True
    if italic: r.font.italic = True
    return r


def _styled(para, text, sz, default_clr=None):
    """Append runs with **bold** / *italic* / `code` colouring."""
    dc = default_clr or C_FG
    for txt, sty in parse_inline(text):
        if   sty == "b": _r(para, txt, FB, sz, C_ACCENT, bold=True)
        elif sty == "i": _r(para, txt, FB, sz, C_PURPLE)
        elif sty == "c": _r(para, txt, FC, sz, C_ACCENT)
        else:            _r(para, txt, FB, sz, dc)


# ═══════════════════════════════════════════════════════════════
#  CHROME  (background, stripe, bar, footer, title bar)
# ═══════════════════════════════════════════════════════════════

def _no_outline(shape):
    shape.line.fill.background()


def chrome(slide, idx, total):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C_BG

    # green left stripe
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, STRIPE, SH)
    s.fill.solid(); s.fill.fore_color.rgb = C_ACCENT; _no_outline(s)

    # progress bar
    bh = Emu(30000)
    bw = int(SW * (idx + 1) / total)
    b = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, SH - bh, bw, bh)
    b.fill.solid(); b.fill.fore_color.rgb = C_ACCENT; _no_outline(b)


def add_footer(slide, text=FOOTER_TEXT):
    tb = slide.shapes.add_textbox(MX, Inches(6.88), CW, Inches(0.35))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    _r(p, "// ", FC, Z["footer"], C_ACCENT)
    _r(p, text, FC, Z["footer"], C_GRAY)


def add_title_bar(slide, title_text):
    """## Title  +  bottom border line."""
    tb = slide.shapes.add_textbox(MX, TITLE_Y, CW, TITLE_H)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    _r(p, "## ", FC, Z["h2"], C_ACCENT, bold=True)
    _r(p, title_text, FC, Z["h2"], C_HEAD, bold=True)

    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MX, BORDER_Y, CW, Emu(19000))
    ln.fill.solid(); ln.fill.fore_color.rgb = C_BORDER; _no_outline(ln)


# ═══════════════════════════════════════════════════════════════
#  BLOCK  →  TEXT-FRAME  RENDERER
# ═══════════════════════════════════════════════════════════════

def _est_inches(blocks):
    h = 0.0
    for bl in blocks:
        t = bl["type"]
        if   t == "h3":    h += 0.42
        elif t == "list":  h += len(bl["items"]) * 0.30
        elif t == "para":  h += max(1, len(bl["text"]) // 100) * 0.30
        elif t == "table": h += len(bl["rows"]) * 0.38
        elif t == "code":  h += max(1, bl["text"].count("\n") + 1) * 0.22
        elif t == "quote": h += max(1, len(bl["text"]) // 100) * 0.35
    return h


def _render_text(tf, blocks, sz=None):
    """Render non-table blocks into an existing TextFrame."""
    sz = sz or Z["body"]
    first_para = True
    first_block = True

    for bl in blocks:
        gap = Pt(10) if not first_block else Pt(0)

        if bl["type"] == "h3":
            p = tf.paragraphs[0] if first_para else tf.add_paragraph()
            p.space_before = Pt(14) if not first_block else Pt(0)
            p.space_after  = Pt(3)
            _r(p, "### ", FC, Z["h3"], C_ACCENT)
            _r(p, bl["text"], FC, Z["h3"], C_FG, bold=True)
            first_para = False

        elif bl["type"] == "list":
            for li, item in enumerate(bl["items"]):
                p = tf.paragraphs[0] if first_para else tf.add_paragraph()
                p.space_before = gap if li == 0 else Pt(2)
                p.space_after  = Pt(3)
                pad = "     " * item["level"]
                bullet = f"{item['num']}. " if item.get("num") else "\u2022  "
                _r(p, pad + bullet, FB, sz, C_ACCENT)
                _styled(p, item["text"], sz)
                first_para = False

        elif bl["type"] == "para":
            p = tf.paragraphs[0] if first_para else tf.add_paragraph()
            p.space_before = gap
            p.space_after  = Pt(4)
            _styled(p, bl["text"], sz)
            first_para = False

        elif bl["type"] == "code":
            for ci, cline in enumerate(bl["text"].split("\n")):
                p = tf.paragraphs[0] if first_para else tf.add_paragraph()
                p.space_before = gap if ci == 0 else Pt(0)
                p.space_after  = Pt(0)
                _r(p, cline or " ", FC, Z["code"], C_ACCENT)
                first_para = False

        elif bl["type"] == "quote":
            p = tf.paragraphs[0] if first_para else tf.add_paragraph()
            p.space_before = gap
            p.space_after  = Pt(4)
            _r(p, "\u2502 ", FC, sz, C_HEAD)
            _styled(p, bl["text"], sz, default_clr=C_GRAY)
            first_para = False

        first_block = False


# ═══════════════════════════════════════════════════════════════
#  NATIVE  TABLE  RENDERER
# ═══════════════════════════════════════════════════════════════

def _table_shape(slide, rows, x, y, w):
    """Create + style a PowerPoint table.  Returns height (EMU)."""
    nr = len(rows)
    nc = max(len(r) for r in rows) if rows else 1
    rh = Inches(0.35)
    h  = rh * nr

    shape = slide.shapes.add_table(nr, nc, x, y, w, h)
    tbl   = shape.table

    # kill default Office table style
    tbl_xml = shape._element.find(".//" + qn("a:tbl"))
    tblPr = tbl_xml.find(qn("a:tblPr"))
    if tblPr is not None:
        for a in list(tblPr.attrib):
            del tblPr.attrib[a]
        tblPr.set("bandRow", "0")
        tblPr.set("bandCol", "0")
        tblPr.set("firstRow", "0")
        for child in list(tblPr):
            tblPr.remove(child)

    col_w = w // nc
    for ci in range(nc):
        tbl.columns[ci].width = col_w

    for ri, row in enumerate(rows):
        for ci in range(nc):
            txt = row[ci] if ci < len(row) else ""
            cell = tbl.cell(ri, ci)

            cell.fill.solid()
            cell.fill.fore_color.rgb = C_CODE_BG if ri == 0 else C_BG

            # borders
            tc = cell._tc
            tcPr = tc.find(qn("a:tcPr"))
            if tcPr is None:
                tcPr = SubElement(tc, qn("a:tcPr"))
            for side in ("lnL", "lnR", "lnT", "lnB"):
                old = tcPr.find(qn(f"a:{side}"))
                if old is not None:
                    tcPr.remove(old)
                ln = SubElement(tcPr, qn(f"a:{side}"))
                ln.set("w", "12700" if (ri == 0 and side == "lnB") else "6350")
                sf = SubElement(ln, qn("a:solidFill"))
                clr_el = SubElement(sf, qn("a:srgbClr"))
                clr_el.set("val", "30363D")

            # cell text
            cell.text = ""
            cell.margin_left  = Emu(72000)
            cell.margin_right = Emu(72000)
            cell.margin_top   = Emu(36000)
            cell.margin_bottom= Emu(36000)
            p = cell.text_frame.paragraphs[0]

            for txt_s, sty in parse_inline(txt):
                r = p.add_run()
                r.text = txt_s
                r.font.name = FB
                if ri == 0:
                    r.font.size = Z["th"]
                    r.font.color.rgb = C_HEAD
                    r.font.bold = True
                elif sty == "b":
                    r.font.size = Z["td"]
                    r.font.color.rgb = C_ACCENT
                    r.font.bold = True
                elif sty == "i":
                    r.font.size = Z["td"]
                    r.font.color.rgb = C_PURPLE
                else:
                    r.font.size = Z["td"]
                    r.font.color.rgb = C_GRAY
    return h


# ═══════════════════════════════════════════════════════════════
#  SLIDE  RENDERERS
# ═══════════════════════════════════════════════════════════════

def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def render_title(prs, d, idx, n):
    slide = _blank(prs)
    chrome(slide, idx, n)

    tb = slide.shapes.add_textbox(MX, Inches(1.9), CW, Inches(2.0))
    tf = tb.text_frame; tf.word_wrap = True
    for j, line in enumerate(d["title"].split("\n")):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        _r(p, line, FC, Z["h1"], C_HEAD, bold=True)
        p.space_after = Pt(6)

    if d["body"]:
        bb = slide.shapes.add_textbox(MX, Inches(4.05), CW, Inches(2.8))
        bf = bb.text_frame; bf.word_wrap = True
        first = True
        for line in d["body"].split("\n"):
            line = line.strip()
            if not line:
                continue
            p = bf.paragraphs[0] if first else bf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            p.space_after = Pt(8)
            _styled(p, line, Z["lead_body"])
            first = False


def render_section(prs, d, idx, n):
    slide = _blank(prs)
    chrome(slide, idx, n)

    tb = slide.shapes.add_textbox(MX, Inches(2.6), CW, Inches(1.2))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _r(p, d["title"], FC, Z["sec_title"], C_HEAD, bold=True)

    body_text = d["body"].strip().split("\n")[0] if d["body"] else ""
    if body_text:
        sb = slide.shapes.add_textbox(MX, Inches(3.9), CW, Inches(0.8))
        sf = sb.text_frame; sf.word_wrap = True
        p2 = sf.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        _r(p2, body_text, FC, Z["sec_sub"], C_GRAY)


def render_lead(prs, d, idx, n):
    slide = _blank(prs)
    chrome(slide, idx, n)

    tb = slide.shapes.add_textbox(MX, Inches(1.8), CW, Inches(1.2))
    tf = tb.text_frame; tf.word_wrap = True
    for j, line in enumerate(d["title"].split("\n")):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        _r(p, line, FC, Z["h1"], C_HEAD, bold=True)
        p.space_after = Pt(6)

    if d["body"]:
        bb = slide.shapes.add_textbox(MX, Inches(3.3), CW, Inches(3.2))
        bf = bb.text_frame; bf.word_wrap = True
        first = True
        for line in d["body"].split("\n"):
            line = line.strip()
            if not line:
                continue
            p = bf.paragraphs[0] if first else bf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            p.space_after = Pt(8)
            _styled(p, line, Z["lead_body"])
            first = False


def render_content(prs, d, idx, n):
    slide = _blank(prs)
    chrome(slide, idx, n)
    add_footer(slide)
    add_title_bar(slide, d["title"])

    blocks = d["blocks"]
    if not blocks:
        return

    # split into (text_group | table) segments
    groups = []
    txt_buf = []
    for bl in blocks:
        if bl["type"] == "table":
            if txt_buf:
                groups.append(("text", txt_buf))
                txt_buf = []
            groups.append(("table", bl))
        else:
            txt_buf.append(bl)
    if txt_buf:
        groups.append(("text", txt_buf))

    # height estimates (inches)
    heights = []
    for gt, gd in groups:
        if gt == "text":
            heights.append(_est_inches(gd))
        else:
            heights.append(len(gd["rows"]) * 0.38)

    avail = BODY_MAX / 914400          # EMU → inches
    total_h = sum(heights) or 1.0
    scale = min(1.0, avail / total_h)

    y = BODY_Y
    for gi, (gt, gd) in enumerate(groups):
        h_in  = heights[gi] * scale
        h_emu = max(int(h_in * 914400), Emu(200000))

        if gt == "text":
            tb = slide.shapes.add_textbox(MX, y, CW, h_emu)
            tf = tb.text_frame
            tf.word_wrap = True
            if scale < 0.88:
                tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            _render_text(tf, gd)
            y += h_emu
        else:
            th = _table_shape(slide, gd["rows"], MX, y, CW)
            y += th + Emu(72000)


# ═══════════════════════════════════════════════════════════════
#  MAIN
# ═══════════════════════════════════════════════════════════════

def main():
    src    = DIR / "slides.md"
    output = DIR / "slides.pptx"

    slides = parse_slides(src)
    print(f"Parsed {len(slides)} slides from {src.name}")

    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH
    total = len(slides)

    for idx, sd in enumerate(slides):
        t = sd["type"]
        if   t == "title":   render_title(prs, sd, idx, total)
        elif t == "section": render_section(prs, sd, idx, total)
        elif t == "lead":    render_lead(prs, sd, idx, total)
        else:                render_content(prs, sd, idx, total)

    prs.save(str(output))
    print(f"Saved {output}  ({total} slides)")


if __name__ == "__main__":
    main()
